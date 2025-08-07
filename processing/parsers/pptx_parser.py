from pptx import Presentation
from typing import Any, Dict, List


class PPTXParser:
    def __init__(self, config: Dict[str, Any]) -> None:
        self.config = config

    def parse(self, filepath: str) -> List[str]:
        prs = Presentation(filepath)
        text_runs: List[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)  # type: ignore
        return text_runs
