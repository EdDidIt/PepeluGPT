from typing import Dict, Any
from pptx import Presentation  # type: ignore

def parse(filepath: str) -> Dict[str, Any]:
    try:
        # Load the PowerPoint presentation
        prs = Presentation(filepath)
        content = ""  # Start as plain string
        
        # Extract text from all slides
        for slide_num, slide in enumerate(prs.slides, 1):
            content += f"=== Slide {slide_num} ===\n"
            
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    shape_text = getattr(shape, "text", "")
                    if isinstance(shape_text, str) and shape_text.strip():
                        content += shape_text + "\n"
            content += "\n"
        
        return {
            "status": "success",
            "filename": filepath,
            "content": content.strip(),
            "error": None
        }
    except Exception as e:
        return {
            "status": "error",
            "filename": filepath,
            "content": None,
            "error": str(e)
        }